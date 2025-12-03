import os
import zipfile
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle
)
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_LEFT
from reportlab.lib import colors

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

os.makedirs("output", exist_ok=True)
os.makedirs("charts", exist_ok=True)

# ============================
#         PDF BUILDER
# ============================

def create_pdf(metrics, chart_paths, insights_text, base_name, anomalies):
    pdf_path = os.path.join("output", f"{base_name}.pdf")

    doc = SimpleDocTemplate(
        pdf_path, pagesize=A4,
        rightMargin=40, leftMargin=40,
        topMargin=50, bottomMargin=40
    )

    styles = getSampleStyleSheet()

    title_style = ParagraphStyle(
        "title", parent=styles["Title"],
        fontName="Helvetica-Bold", fontSize=26,
        leading=30, textColor=colors.HexColor("#0B1D4A"),
        spaceAfter=20
    )

    section_style = ParagraphStyle(
        "section", parent=styles["Heading2"],
        fontName="Helvetica-Bold", fontSize=16,
        textColor=colors.HexColor("#0B1D4A"),
        spaceBefore=20, spaceAfter=10
    )

    body = ParagraphStyle(
        "body", parent=styles["BodyText"],
        fontName="Helvetica", fontSize=11,
        leading=17, alignment=TA_LEFT,
        spaceAfter=6
    )

    story = []

    # Title
    story.append(Paragraph("Automated Insight Engine", title_style))
    story.append(Spacer(1, 10))

    # EXEC SUMMARY
    story.append(Paragraph("EXECUTIVE SUMMARY", section_style))

    for line in insights_text.split("\n"):
        if line.strip():
            clean = "• " + line.lstrip("*•- ").strip()
            story.append(Paragraph(clean, body))

    story.append(Spacer(1, 10))

    # ANOMALIES
    story.append(Paragraph("DETECTED ANOMALIES", section_style))

    if anomalies:
        for a in anomalies[:7]:
            story.append(Paragraph("• " + str(a), body))
    else:
        story.append(Paragraph("• No major anomalies detected.", body))

    story.append(Spacer(1, 10))

    # METRICS
    story.append(Paragraph("KEY METRICS", section_style))

    rows = [["Metric", "Value"]]
    for k, v in metrics.items():
        if isinstance(v, dict):
            val = "<br/>".join([f"{kk}: {vv}" for kk, vv in v.items()])
            rows.append([k, Paragraph(val, body)])
        else:
            rows.append([k, str(v)])

    table = Table(rows, colWidths=[150, 330])
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#E1E6F0")),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("GRID", (0, 0), (-1, -1), 0.3, colors.HexColor("#777777")),
        ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
        ("VALIGN", (0, 0), (-1, -1), "TOP")
    ]))

    story.append(table)
    story.append(Spacer(1, 20))

    # CHARTS
    story.append(Paragraph("VISUAL INSIGHTS", section_style))

    for ch in chart_paths:
        story.append(Image(ch, width=460, height=240))
        story.append(Spacer(1, 12))

    doc.build(story)
    return pdf_path

# ============================
#     PPT HELPERS
# ============================

def add_bullets(slide, lines, top=1.2, base_size=22):
    """
    Smart keynote-style bullets:
    • Auto font shrink based on item count
    • Overflow-safe
    • Clean margins + spacing
    """

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(12), Inches(5.6))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.clear()

    # Auto scale font based on bullet count
    n = len(lines)

    if n > 10:
        size = 16
    elif n > 7:
        size = 18
    elif n > 5:
        size = 20
    else:
        size = base_size

    for line in lines:
        clean = line.strip().lstrip("*•- ").strip()
        if not clean:
            continue

        p = tf.add_paragraph()
        p.text = clean
        p.level = 1
        p.font.size = Pt(size)
        p.font.name = "Helvetica"
        p.font.color.rgb = RGBColor(40, 40, 40)


# ==========================================
#               PPT BUILDER
# ==========================================
def create_ppt(metrics, charts, insights, base_name, anomalies):

    prs = Presentation()
    prs.slide_width = Inches(13.3)
    prs.slide_height = Inches(7.5)

    # --- TITLE SLIDE ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = tb.text_frame
    tf.text = "Automated Insight Engine"
    tf.paragraphs[0].font.size = Pt(52)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # --- EXEC SUMMARY SLIDE ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tx = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(10), Inches(1))
    tx.text_frame.text = "Executive Summary"
    tx.text_frame.paragraphs[0].font.size = Pt(40)

    summary_lines = [x for x in insights.split("\n") if x.strip()]
    summary_lines = summary_lines[:10]  # hard cap

    add_bullets(slide, summary_lines, top=1.4, base_size=22)

    # --- METRICS SLIDE ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    mt = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(10), Inches(1))
    mt.text_frame.text = "Key Metrics"
    mt.text_frame.paragraphs[0].font.size = Pt(40)

    table = slide.shapes.add_table(
        len(metrics) + 1, 2, Inches(0.7), Inches(1.4), Inches(11), Inches(4)
    ).table

    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"

    i = 1
    for k, v in metrics.items():
        table.cell(i, 0).text = str(k)
        table.cell(i, 1).text = str(v)
        i += 1

    for r in range(len(metrics) + 1):
        for c in range(2):
            for p in table.cell(r, c).text_frame.paragraphs:
                p.font.size = Pt(20)

    # --- ANOMALIES SLIDE ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    at = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(10), Inches(1))
    at.text_frame.text = "Detected Anomalies"
    at.text_frame.paragraphs[0].font.size = Pt(40)

    # Normalize anomalies
    anomaly_lines = []

    if anomalies:
        for a in anomalies[:7]:
            if isinstance(a, dict):
                anomaly_lines.append(
                    f"{a.get('date','?')} – {a.get('city','?')}: CTR shift {a.get('ctr_change',0):.2f}x"
                )
            else:
                anomaly_lines.append(str(a))
    else:
        anomaly_lines.append("No anomalies detected.")

    add_bullets(slide, anomaly_lines, top=1.4, base_size=22)

    # --- CHART SLIDES ---
    if isinstance(charts, dict):
        items = charts.items()
    else:
        items = [(os.path.basename(ch), ch) for ch in charts]

    for label, path in items:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(10), Inches(1))
        tb.text_frame.text = label.replace("_", " ").title()
        tb.text_frame.paragraphs[0].font.size = Pt(36)

        slide.shapes.add_picture(path, Inches(1), Inches(1.3), width=Inches(11))

    ppt_path = f"output/{base_name}_report.pptx"
    prs.save(ppt_path)
    return ppt_path

# ============================
# ZIP BUILDER
# ============================

def create_zip(pdf_path, ppt_path, chart_paths, base_name):
    zip_path = f"output/{base_name}.zip"
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as z:
        z.write(pdf_path, os.path.basename(pdf_path))
        z.write(ppt_path, os.path.basename(ppt_path))
        for ch in chart_paths:
            z.write(ch, os.path.basename(ch))
    return zip_path